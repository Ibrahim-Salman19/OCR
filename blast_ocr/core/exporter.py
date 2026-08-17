"""
blast_ocr.core.exporter

Multi-format Document Serializer and Exporter.
Transforms OCR extractions into Markdown with Frontmatter, styled DOCX (with native tables),
EPUB 3.0, Plain Text, Searchable Sandwich PDF, and JSON manifests.
"""

import os
import re
import sys
import json
import logging
from pathlib import Path
from typing import Optional, Any, List, Dict
from pptx import Presentation as DefaultPresentation
from docx import Document as DefaultDocument

import defusedxml
defusedxml.defuse_stdlib()

from blast_ocr.core.exceptions import OCREngineError
from blast_ocr.core.models import ExportBundle

logger = logging.getLogger(__name__)


def _get_presentation_cls():
    extractor_mod = sys.modules.get("blast_ocr.core.extractor")
    if extractor_mod and hasattr(extractor_mod, "Presentation"):
        return extractor_mod.Presentation
    return DefaultPresentation


def _get_document_cls():
    extractor_mod = sys.modules.get("blast_ocr.core.extractor")
    if extractor_mod and hasattr(extractor_mod, "Document"):
        return extractor_mod.Document
    return DefaultDocument


def sanitize_for_xml(text: Optional[str]) -> str:
    """Removes control characters that are invalid in XML / DOCX schemas."""
    if not text:
        return ""
    return re.sub(r"[^\x09\x0A\x0D\x20-\uD7FF\uE000-\uFFFD\u10000-\u10FFFF]", "", text)


def _render_docx_content(doc: Any, text: str) -> None:
    """Renders text lines, headings, page breaks, and Markdown tables into native Word doc."""
    lines = text.split("\n")
    idx = 0
    while idx < len(lines):
        line = lines[idx].strip()
        
        # Check if line is part of a markdown table (e.g. '| col1 | col2 |')
        if line.startswith("|") and line.endswith("|") and "|" in line[1:-1]:
            table_lines = []
            while idx < len(lines) and lines[idx].strip().startswith("|") and lines[idx].strip().endswith("|"):
                tl = lines[idx].strip()
                # Skip separator rows like '|---|---|'
                if not re.match(r"^\|[\s\-:|]+\|$", tl):
                    cells = [c.strip() for c in tl.strip("|").split("|")]
                    table_lines.append(cells)
                idx += 1
            
            if table_lines:
                num_cols = max(len(r) for r in table_lines)
                table = doc.add_table(rows=len(table_lines), cols=num_cols)
                table.style = "Table Grid"
                for r_i, row_data in enumerate(table_lines):
                    for c_i, cell_text in enumerate(row_data):
                        if c_i < num_cols:
                            table.cell(r_i, c_i).text = cell_text
                continue
        elif line.startswith("# "):
            doc.add_heading(line.replace("# ", ""), level=1)
        elif line.startswith("## "):
            doc.add_heading(line.replace("## ", ""), level=2)
        elif line.startswith("### "):
            doc.add_heading(line.replace("### ", ""), level=3)
        elif line.startswith("---"):
            doc.add_page_break()
        else:
            if line:
                doc.add_paragraph(line)
        idx += 1


def save_output(
    text: str,
    base_name: str,
    output_dir: str,
    doc_model: Optional[Any] = None,
    page_images: Optional[List[Any]] = None,
    page_results: Optional[List[Dict[str, Any]]] = None,
    metadata: Optional[Dict[str, Any]] = None,
) -> ExportBundle:
    """
    Saves output to Markdown, DOCX, TXT, EPUB, Searchable PDF, and JSON manifest.
    
    Returns:
        ExportBundle instance.
    """
    os.makedirs(output_dir, exist_ok=True)
    clean_text = sanitize_for_xml(text)

    # 1. Save Markdown (with Frontmatter if metadata present)
    md_path = Path(output_dir) / f"{base_name}.md"
    md_content = clean_text
    if metadata and isinstance(metadata, dict):
        fm_lines = ["---"]
        for k, v in sorted(metadata.items()):
            fm_lines.append(f"{k}: {v}")
        fm_lines.append("---")
        fm_lines.append("")
        md_content = "\n".join(fm_lines) + clean_text

    with open(md_path, "w", encoding="utf-8") as f:
        f.write(md_content)

    # 2. Save Plain Text
    txt_path = Path(output_dir) / f"{base_name}.txt"
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(clean_text)

    # 3. Save DOCX
    docx_path = Path(output_dir) / f"{base_name}.docx"
    try:
        DocCls = _get_document_cls()
        doc = DocCls()
        doc.add_heading(base_name, 0)
        _render_docx_content(doc, clean_text)
        doc.save(docx_path)
    except Exception as e:
        logger.error(f"DOCX generation failed: {e}")
        docx_path = None

    # 4. Save EPUB if Document model present
    epub_path = None
    if doc_model is not None:
        try:
            epub_path = Path(output_dir) / f"{base_name}.epub"
            from blast_ocr.core.book_intelligence import BookProcessor
            BookProcessor.export_epub(doc_model, str(epub_path), title=base_name)
        except Exception as e:
            logger.warning(f"EPUB generation failed: {e}")
            epub_path = None

    # 5. Save Searchable PDF if page images & OCR results are provided
    pdf_path = None
    if page_images and page_results:
        try:
            from blast_ocr.core.searchable_pdf import SearchablePDFGenerator
            out_pdf = str(Path(output_dir) / f"{base_name}.pdf")
            SearchablePDFGenerator.create_from_page_images(
                page_images=page_images,
                page_ocr_results=page_results,
                output_pdf_path=out_pdf,
                title=base_name,
            )
            pdf_path = Path(out_pdf)
        except Exception as e:
            logger.warning(f"Searchable PDF generation failed: {e}")
            pdf_path = None

    # 6. Save JSON Manifest
    manifest_path = Path(output_dir) / f"{base_name}_manifest.json"
    manifest_data = {
        "document_name": base_name,
        "metadata": metadata or {},
        "generated_files": {
            "markdown": str(md_path),
            "docx": str(docx_path) if docx_path else None,
            "txt": str(txt_path),
            "epub": str(epub_path) if epub_path else None,
            "pdf": str(pdf_path) if pdf_path else None,
        },
    }
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest_data, f, indent=2)

    return ExportBundle(
        markdown_path=md_path,
        docx_path=docx_path,
        txt_path=txt_path,
        epub_path=epub_path,
        pdf_path=pdf_path,
        manifest_path=manifest_path,
    )


def extract_from_pptx(pptx_path: str) -> str:
    """Extracts text from slides, including notes and tables."""
    text_content = []
    try:
        PrsCls = _get_presentation_cls()
        prs = PrsCls(pptx_path)
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            slide_text.append(f"## Slide {i}")

            # 1. Shapes Text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

                # 2. Tables
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            [cell.text_frame.text for cell in row.cells]
                        )
                        slide_text.append(f"| {row_text} |")

            # 3. Notes
            if (
                hasattr(slide, "has_notes_slide")
                and slide.has_notes_slide
                and slide.notes_slide.notes_text_frame
            ):
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    slide_text.append(f"> **Notes:** {notes}")

            text_content.append("\n".join(slide_text))

        return "\n\n---\n\n".join(text_content)
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        raise OCREngineError(f"PPTX extraction failed: {e}") from e
