"""
Layer 3 Tool: Text Extractor
Phase: Architect

Handles deterministic extraction from:
1. PPTX (Native text + Table text)
2. Images (EasyOCR with preprocessing)
3. PDF (Conversion to images -> EasyOCR)
"""

import os
import sys
import tempfile
import shutil
import cv2
import numpy as np
import easyocr
import glob
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image
from docx import Document

# Initialize EasyOCR Reader once (Global)
try:
    READER = easyocr.Reader(['en'], gpu=False)
except Exception as e:
    print(f"[!] EasyOCR Init Failed: {e}")
    READER = None

def preprocess_image(image_path, target_width=2000):
    """
    Standard B.L.A.S.T. Preprocessing:
    Gray -> Denoise -> Deskew -> Resize -> Adaptive Thresh
    """
    try:
        # Load
        img = cv2.imdecode(np.fromfile(image_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img is None:
            return None

        # 1. Gray
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        
        # 2. Denoise (Fast)
        gray = cv2.fastNlMeansDenoising(gray, None, h=10, templateWindowSize=7, searchWindowSize=21)

        # 3. Deskew (Simple)
        coords = np.column_stack(np.where(gray > 0))
        angle = cv2.minAreaRect(coords)[-1]
        if angle < -45: angle = -(90 + angle)
        else: angle = -angle
        if abs(angle) > 0.2:
            (h, w) = gray.shape
            center = (w // 2, h // 2)
            M = cv2.getRotationMatrix2D(center, angle, 1.0)
            gray = cv2.warpAffine(gray, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)

        # 4. Resize
        h, w = gray.shape
        if w < target_width:
            scale = target_width / float(w)
            gray = cv2.resize(gray, (int(w * scale), int(h * scale)), interpolation=cv2.INTER_CUBIC)

        # 5. Adaptive Threshold
        bin_img = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                        cv2.THRESH_BINARY, 21, 10)
        return bin_img
    except Exception as e:
        print(f"[!] Preprocessing failed for {image_path}: {e}")
        return None

def extract_from_image(image_path):
    """Returns text from an image using EasyOCR."""
    if READER is None:
        return "[ERROR: EasyOCR not initialized]"
    
    # Preprocess
    processed = preprocess_image(image_path)
    if processed is None:
        return "[ERROR: Image load failed]"
    
    # EasyOCR expects RGB or File
    # Convert binary numpy back to valid input for EasyOCR (or pass directly)
    # EasyOCR handles numpy arrays
    try:
        results = READER.readtext(processed, detail=0, paragraph=True)
        return "\n\n".join(results)
    except Exception as e:
        return f"[ERROR: OCR failed: {e}]"

def extract_from_pptx(pptx_path):
    """Extracts text from slides, including notes and tables."""
    text_content = []
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            slide_text.append(f"## Slide {i}")
            
            # 1. Shapes Text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                
                # 2. Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text_frame.text for cell in row.cells])
                        slide_text.append(f"| {row_text} |")
            
            # 3. Notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    slide_text.append(f"> **Notes:** {notes}")
            
            text_content.append("\n".join(slide_text))
            
        return "\n\n---\n\n".join(text_content)
    except Exception as e:
        return f"[ERROR: PPTX extraction failed: {e}]"

def extract_from_pdf(pdf_path, poppler_path=None):
    """Converts PDF to images then runs OCR."""
    text_content = []
    try:
        # Create temp dir for pages
        with tempfile.TemporaryDirectory() as temp_dir:
            # Convert
            try:
                pages = convert_from_path(pdf_path, dpi=300, output_folder=temp_dir, poppler_path=poppler_path)
            except Exception as e:
                return f"[ERROR: PDF to Image conversion failed. Is Poppler installed? Details: {e}]"
            
            # Iterate
            image_files = sorted(glob.glob(os.path.join(temp_dir, "*.ppm"))) # pdf2image saves as ppm often
            if not image_files:
                # Fallback if saved differently or list is empty, retry iterating generated PIL objects if needed
                # But convert_from_path returns PIL images if output_folder is None.
                # Let's use the PIL list directly.
                pass

            for i, page_img in enumerate(pages, start=1):
                # Save as PNG for preprocessing
                page_path = os.path.join(temp_dir, f"page_{i}.png")
                page_img.save(page_path, "PNG")
                
                # OCR
                page_text = extract_from_image(page_path)
                text_content.append(f"## Page {i}\n\n{page_text}")
                
        return "\n\n---\n\n".join(text_content)
            
    except Exception as e:
        return f"[ERROR: PDF processing crashed: {e}]"

def save_output(text, base_name, output_dir):
    """Saves to Markdown and DOCX."""
    os.makedirs(output_dir, exist_ok=True)
    
    # MD
    md_path = os.path.join(output_dir, f"{base_name}.md")
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(text)
        
    # DOCX
    docx_path = os.path.join(output_dir, f"{base_name}.docx")
    doc = Document()
    doc.add_heading(base_name, 0)
    for line in text.split('\n'):
        line = line.strip()
        if line.startswith('## '):
            doc.add_heading(line.replace('## ', ''), level=2)
        elif line.startswith('---'):
            doc.add_page_break()
        else:
            if line:
                doc.add_paragraph(line)
    doc.save(docx_path)
    
    return md_path, docx_path

if __name__ == "__main__":
    # Simple CLI for testing this module directly
    if len(sys.argv) < 2:
        print("Usage: python text_extractor.py <file_path>")
        sys.exit(1)
        
    fpath = sys.argv[1]
    print(f"[-] Processing: {fpath}")
    
    if fpath.endswith(".pptx"):
        res = extract_from_pptx(fpath)
    elif fpath.endswith(".pdf"):
        res = extract_from_pdf(fpath)
    else:
        res = extract_from_image(fpath)
        
    print(f"[-] Extracted {len(res)} characters.")
    save_output(res, "test_output", ".")
    print("[+] Done.")
