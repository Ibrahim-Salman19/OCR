import os
import sys
# Add project root to path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from blast_ocr.core.extractor import extract_from_pptx

def create_test_pptx(filename="test_verify.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello PPTX World"
    subtitle.text = "This is a subtitle"
    
    # Add a table slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide2.shapes
    rows, cols = 2, 2
    left = top = width = height = 100000
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Row1 Col1"
    table.cell(0, 1).text = "Row1 Col2"
    table.cell(1, 0).text = "Row2 Col1"
    table.cell(1, 1).text = "Row2 Col2"
    
    prs.save(filename)
    return filename

def verify():
    pptx_path = create_test_pptx()
    print(f"[-] Created {pptx_path}")
    
    print("[-] Extracting text...")
    text = extract_from_pptx(pptx_path)
    print(f"[-] Extracted:\n{text}")
    
    if "Hello PPTX World" in text and "Row1 Col1" in text:
        print("[+] PPTX verification PASSED")
    else:
        print("[!] PPTX verification FAILED")
        
    if os.path.exists(pptx_path):
        os.remove(pptx_path)

if __name__ == "__main__":
    verify()
